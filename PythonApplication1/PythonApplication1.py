# First, run 'pip install python-pptx' in your terminal
from pptx import Presentation

prs = Presentation()

# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Beyond the Hype: Applied AI & Real-World Application"
slide.placeholders[1].text = "Navigating the AI Era as Learners, Innovators, and Leaders"

# Slide 2: Shift to Applied AI
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "The Shift to Applied AI"
slide.placeholders[1].text = "• The Past: Theoretical AI (1950–2022) - Hidden in complex labs\n• The Present: Applied AI Utility Layer - Works quietly like electricity\n• The Great Equalizer: Clear articulation is your new programming language"

# Slide 3: Education Transformation
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "The Educational Transformation"
slide.placeholders[1].text = "• End of Rote Memorization: Focus shifts to connecting dots\n• Hyper-Personalized Tutors: Material adapts to your personal interests\n• Intellectual Sparring Partner: Moves from learning 'what' to 'how' to think"

# Save the presentation
prs.save("Practical_AI_Presentation.pptx")
print("Your PowerPoint file 'Practical_AI_Presentation.pptx' has been successfully created!")
